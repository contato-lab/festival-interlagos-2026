#!/usr/bin/env python3
"""
generate-meta-creatives-pptx.py

Gera uma apresentação PowerPoint (.pptx) com os top criativos do Meta Ads
do Festival Interlagos 2026, baseado no arquivo analise-meta-criativos.json.

Cada criativo ganha um slide com:
- Thumbnail (baixada do Meta CDN)
- Métricas (compras, gasto, CPP, CTR, impressões, cliques)
- Nome do anúncio + adset + campanha
- Link direto pro Ads Manager

Output: top-criativos-meta.pptx
"""

import json
import os
import io
import urllib.request
import urllib.parse
import sys
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ─── Paleta (identidade Lime) ───────────────────────────────────────────
COR_BG_ESCURO  = RGBColor(0x0A, 0x0A, 0x0A)
COR_LIME       = RGBColor(0xC8, 0xFF, 0x00)
COR_TEXT       = RGBColor(0xFF, 0xFF, 0xFF)
COR_TEXT_DIM   = RGBColor(0xAA, 0xAA, 0xAA)
COR_ACCENT_BLUE  = RGBColor(0x5D, 0xAE, 0xFF)
COR_ACCENT_GREEN = RGBColor(0x4D, 0xFF, 0x91)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)


def baixa_thumb(url, dest_path):
    if not url:
        return False
    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=20) as r:
            data = r.read()
        with open(dest_path, 'wb') as f:
            f.write(data)
        return True
    except Exception as e:
        print(f'  ! falha baixando {url[:80]}...: {e}')
        return False


def add_bg(slide, color=COR_BG_ESCURO):
    """Pinta o fundo do slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text,
             font_size=14, bold=False, color=COR_TEXT, align=PP_ALIGN.LEFT,
             font_name='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def fmt_brl(v):
    if v is None: return '—'
    return 'R$ ' + f'{v:,.2f}'.replace(',', 'X').replace('.', ',').replace('X', '.')


def slide_capa(prs, periodo):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)

    # Lime brand bar (faixa lime no topo)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = COR_LIME
    bar.line.fill.background()

    # Logo / Tag
    add_text(slide, Inches(0.6), Inches(0.5), Inches(4), Inches(0.4),
             'LIME · AGÊNCIA', font_size=11, bold=True, color=COR_LIME,
             font_name='Arial Black')

    # Título grande
    add_text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(1.5),
             'TOP CRIATIVOS', font_size=68, bold=True, color=COR_TEXT,
             font_name='Arial Black')

    add_text(slide, Inches(0.6), Inches(3.5), Inches(12), Inches(0.7),
             'META ADS · FESTIVAL INTERLAGOS 2026',
             font_size=24, bold=True, color=COR_LIME, font_name='Arial')

    # Periodo
    add_text(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.5),
             f'Análise dos últimos {periodo["dias"]} dias  ·  {periodo["since"]} a {periodo["until"]}',
             font_size=16, color=COR_TEXT_DIM, font_name='Calibri')

    # Contas
    add_text(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
             '2 contas: Principal (act_2044706169171045) + Nova Motos (act_1326431289216611)',
             font_size=12, color=COR_TEXT_DIM)

    # Footer
    add_text(slide, Inches(0.6), Inches(6.95), Inches(12), Inches(0.3),
             f'Gerado em {datetime.now().strftime("%d/%m/%Y %H:%M")} · Festival Interlagos 2026',
             font_size=10, color=COR_TEXT_DIM)


def slide_overview_tabela(prs, titulo, subtitulo, ads):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Título
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             titulo, font_size=32, bold=True, color=COR_TEXT, font_name='Arial Black')
    add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
             subtitulo, font_size=14, color=COR_LIME, font_name='Calibri')

    # Cabeçalho da tabela
    cols = [
        ('#',         Inches(0.4)),
        ('Conta',     Inches(1.1)),
        ('Compras',   Inches(0.9)),
        ('Gasto',     Inches(1.2)),
        ('CPP',       Inches(1.1)),
        ('CTR',       Inches(0.7)),
        ('Anúncio',   Inches(7.6)),
    ]
    top = Inches(1.8)
    left_start = Inches(0.6)

    # Header row
    x = left_start
    for label, w in cols:
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, w, Inches(0.45))
        cell.fill.solid(); cell.fill.fore_color.rgb = COR_LIME
        cell.line.fill.background()
        tf = cell.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        run = p.add_run(); run.text = label
        run.font.size = Pt(11); run.font.bold = True
        run.font.color.rgb = COR_BG_ESCURO
        run.font.name = 'Arial'
        x += w

    # Linhas
    row_h = Inches(0.5)
    y = top + Inches(0.5)
    for idx, ad in enumerate(ads, 1):
        x = left_start
        bg_color = RGBColor(0x18, 0x18, 0x18) if idx % 2 == 0 else RGBColor(0x12, 0x12, 0x12)
        for col_label, w in cols:
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, row_h)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg_color
            cell.line.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
            cell.line.width = Pt(0.5)
            tf = cell.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.08)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            if col_label == '#':
                run.text = str(idx); run.font.bold = True; run.font.color.rgb = COR_LIME
            elif col_label == 'Conta':
                run.text = ad['account']
                run.font.color.rgb = COR_ACCENT_BLUE if ad['account'] == 'Nova Motos' else COR_TEXT_DIM
            elif col_label == 'Compras':
                run.text = str(ad['purchases']); run.font.bold = True; run.font.color.rgb = COR_LIME
            elif col_label == 'Gasto':
                run.text = fmt_brl(ad['spend']); run.font.color.rgb = COR_TEXT
            elif col_label == 'CPP':
                run.text = fmt_brl(ad['cpp']); run.font.color.rgb = COR_ACCENT_GREEN
            elif col_label == 'CTR':
                run.text = f"{ad['ctr']}%"; run.font.color.rgb = COR_TEXT
            elif col_label == 'Anúncio':
                run.text = (ad['ad_name'] or '')[:65]; run.font.color.rgb = COR_TEXT
            run.font.size = Pt(10); run.font.name = 'Calibri'
            x += w
        y += row_h


def slide_criativo(prs, rank, ad, badge_label, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Badge top-left
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(2.3), Inches(0.4))
    badge.fill.solid(); badge.fill.fore_color.rgb = COR_LIME
    badge.line.fill.background()
    tf = badge.text_frame; tf.margin_left = Inches(0.1); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'#{rank} de {total} · {badge_label}'
    run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = COR_BG_ESCURO
    run.font.name = 'Arial Black'

    # Nome do anúncio
    nome = (ad['ad_name'] or '(sem nome)').strip()
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
             nome, font_size=28, bold=True, color=COR_TEXT, font_name='Arial Black')

    # Thumbnail (esquerda)
    thumb_path = ad.get('_local_thumb', '')
    if thumb_path and os.path.exists(thumb_path):
        try:
            slide.shapes.add_picture(thumb_path, Inches(0.6), Inches(2.0),
                                     width=Inches(4.5), height=Inches(4.5))
        except Exception as e:
            print(f'   ! falha inserindo imagem {thumb_path}: {e}')
            # Placeholder
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.5))
            ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
            ph.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.0), Inches(4.5), Inches(4.5))
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
        ph.line.color.rgb = RGBColor(0x44, 0x44, 0x44)
        tf = ph.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = '(thumbnail indisponível)'
        run.font.size = Pt(12); run.font.color.rgb = COR_TEXT_DIM

    # Bloco direito: 6 KPIs em grid 2x3
    kpi_data = [
        ('Compras',   str(ad['purchases']),        COR_LIME),
        ('Gasto',     fmt_brl(ad['spend']),        COR_TEXT),
        ('CPP',       fmt_brl(ad['cpp']),          COR_ACCENT_GREEN),
        ('CTR',       f"{ad['ctr']}%",             COR_ACCENT_BLUE),
        ('Impressões',f"{ad['impressions']:,}".replace(',', '.'), COR_TEXT),
        ('Cliques',   f"{ad['clicks']:,}".replace(',', '.'),      COR_TEXT),
    ]
    kpi_left = Inches(5.6)
    kpi_top  = Inches(2.0)
    kpi_w    = Inches(3.7)
    kpi_h    = Inches(1.2)
    gap      = Inches(0.15)

    for i, (label, value, color) in enumerate(kpi_data):
        row = i // 2
        col = i % 2
        x = kpi_left + (kpi_w + gap) * col
        y = kpi_top  + (kpi_h + gap) * row
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, kpi_w, kpi_h)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x16)
        box.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        box.line.width = Pt(0.75)
        # Label
        add_text(slide, x, y + Inches(0.1), kpi_w, Inches(0.3),
                 label.upper(), font_size=10, color=COR_TEXT_DIM,
                 font_name='Arial', align=PP_ALIGN.CENTER)
        # Value
        add_text(slide, x, y + Inches(0.4), kpi_w, Inches(0.7),
                 value, font_size=22, bold=True, color=color,
                 font_name='Arial Black', align=PP_ALIGN.CENTER)

    # Bloco inferior: metadata + link
    meta_top = Inches(6.65)
    add_text(slide, Inches(5.6), meta_top, Inches(7.2), Inches(0.3),
             f"📁 Conta: {ad['account']}  ·  Adset: {(ad['adset_name'] or '')[:60]}",
             font_size=10, color=COR_TEXT_DIM, font_name='Calibri')
    add_text(slide, Inches(5.6), meta_top + Inches(0.25), Inches(7.2), Inches(0.3),
             f"🎯 Campanha: {(ad['campaign_name'] or '')[:80]}",
             font_size=10, color=COR_TEXT_DIM, font_name='Calibri')

    # Link clicável
    account_num = ad['account_id'].replace('act_', '')
    link_url = f"https://business.facebook.com/adsmanager/manage/ads/edit?act={account_num}&selected_ad_ids={ad['ad_id']}"
    tb = slide.shapes.add_textbox(Inches(5.6), meta_top + Inches(0.55), Inches(7.2), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '🔗 Abrir no Meta Ads Manager →'
    run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = COR_LIME
    run.font.name = 'Arial'
    run.hyperlink.address = link_url


def slide_insights(prs, top10, remktg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             'INSIGHTS & RECOMENDAÇÕES', font_size=32, bold=True,
             color=COR_LIME, font_name='Arial Black')

    insights = [
        ('🥇', '#1 GERAL é "VENDAS ABERTAS"',
         '40 compras (2× mais que o segundo). Mantenha rodando — criativo evergreen.'),
        ('💰', 'Melhor CPP da janela: GUIA EDIÇÃO MOTO',
         'R$ 22,61 por compra (Conta Nova Motos). 3× mais barato que a média. Vale escalar.'),
        ('📊', 'Conta Nova Motos é 2,9× mais eficiente',
         'CPP médio R$ 31 vs R$ 65 da Principal. Avaliar migrar mais budget.'),
        ('🔄', '4 dos 10 melhores são REMARKETING',
         'A campanha [RMKT] domina pos #4, #6, #9, #10. Quem visita o site engaja.'),
        ('🚀', 'Maior CTR (2,93%): "VIRADA DE LOTE JARULE"',
         'Combinar headliner + urgência de virada de lote é a fórmula campeã.'),
        ('🎯', 'STREET 108 e D2 aparecem 2× cada',
         'Criativos robustos em múltiplos adsets — replicar pattern em outros temas.'),
    ]

    y = Inches(1.4)
    for icon, titulo, desc in insights:
        # Icon
        add_text(slide, Inches(0.6), y, Inches(0.7), Inches(0.6),
                 icon, font_size=24, color=COR_LIME, align=PP_ALIGN.CENTER)
        # Title
        add_text(slide, Inches(1.4), y, Inches(11.4), Inches(0.4),
                 titulo, font_size=15, bold=True, color=COR_TEXT, font_name='Arial Black')
        # Desc
        add_text(slide, Inches(1.4), y + Inches(0.35), Inches(11.4), Inches(0.4),
                 desc, font_size=12, color=COR_TEXT_DIM, font_name='Calibri')
        y += Inches(0.92)


def main():
    json_path = 'analise-meta-criativos.json'
    with open(json_path, encoding='utf-8') as f:
        data = json.load(f)

    top10 = data['top10_geral']
    remktg = data['top_remarketing_principal']
    periodo = data['periodo']

    # Baixa thumbnails
    os.makedirs('_thumbs', exist_ok=True)
    todos = top10 + remktg
    print(f'Baixando {len(todos)} thumbnails...')
    for ad in todos:
        url = ad.get('thumbnail_url') or ad.get('image_url') or ''
        if not url:
            ad['_local_thumb'] = ''
            continue
        dest = os.path.join('_thumbs', f"{ad['ad_id']}.jpg")
        if os.path.exists(dest) and os.path.getsize(dest) > 100:
            ad['_local_thumb'] = dest
            continue
        ok = baixa_thumb(url, dest)
        ad['_local_thumb'] = dest if ok else ''

    # Build PPTX
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Capa
    slide_capa(prs, periodo)
    # Overview geral
    slide_overview_tabela(
        prs,
        'TOP 10 GERAL · MAIS COMPRAS',
        f'Últimos {periodo["dias"]} dias · 2 contas Meta Ads · ordenado por compras',
        top10,
    )
    # Individuais geral
    for i, ad in enumerate(top10, 1):
        slide_criativo(prs, i, ad, 'TOP GERAL', total=len(top10))
    # Overview remarketing
    slide_overview_tabela(
        prs,
        'TOP REMARKETING · CONTA PRINCIPAL',
        'Campanha [LM] [RMKT] [FESTIVAL INTERLAGOS] [17.04.2026]',
        remktg,
    )
    # Individuais remarketing
    for i, ad in enumerate(remktg, 1):
        slide_criativo(prs, i, ad, 'TOP RMKT', total=len(remktg))
    # Insights
    slide_insights(prs, top10, remktg)

    out = 'top-criativos-meta.pptx'
    prs.save(out)
    print(f'\n✅ Salvo: {out}')
    print(f'   {len(prs.slides)} slides')


if __name__ == '__main__':
    main()
